from __future__ import annotations

from pathlib import Path
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "print_support"
ASSET_DIR = OUT_DIR / "ppt_assets"
PPTX_PATH = Path(os.environ.get("SPECTRA_PPTX_PATH", OUT_DIR / "SPECTRA_PROJECT_PRESENTATION.pptx"))


COLORS = {
    "ink": RGBColor(20, 27, 38),
    "muted": RGBColor(90, 99, 113),
    "blue": RGBColor(37, 99, 235),
    "teal": RGBColor(15, 118, 110),
    "red": RGBColor(220, 38, 38),
    "panel": RGBColor(245, 247, 250),
    "dark": RGBColor(15, 23, 42),
    "white": RGBColor(255, 255, 255),
}


def font(size: int = 28, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
    ]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size)
        except OSError:
            continue
    return ImageFont.load_default()


def make_placeholder(path: Path, title: str, subtitle: str, bg: tuple[int, int, int], accent: tuple[int, int, int]) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), bg)
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((90, 90, w - 90, h - 90), radius=34, outline=accent, width=6)
    draw.rectangle((90, 90, 220, h - 90), fill=accent)
    draw.text((270, 285), title, fill=(255, 255, 255) if sum(bg) < 300 else (20, 27, 38), font=font(72, True))
    draw.text((275, 390), subtitle, fill=(226, 232, 240) if sum(bg) < 300 else (72, 83, 99), font=font(34))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def make_pipeline_image(path: Path) -> None:
    w, h = 1800, 720
    img = Image.new("RGB", (w, h), (247, 249, 252))
    draw = ImageDraw.Draw(img)
    stages = [
        ("Raw images", "phone capture"),
        ("Smart select", "quality frames"),
        ("Keyframes", "normalized input"),
        ("Depth", "Depth Anything V2"),
        ("SfM", "camera poses"),
        ("Fusion", "colored 3D cloud"),
        ("Cleanup", "noise reduction"),
        ("Mesh", "PLY / OBJ"),
    ]
    box_w, box_h = 185, 118
    gap = 32
    x = 55
    y = 245
    for i, (title, subtitle) in enumerate(stages):
        color = (37, 99, 235) if i in {0, 1, 2} else (15, 118, 110) if i in {3, 4, 5} else (124, 58, 237)
        draw.rounded_rectangle((x, y, x + box_w, y + box_h), radius=20, fill=color)
        draw.text((x + 18, y + 25), title, fill=(255, 255, 255), font=font(28, True))
        draw.text((x + 18, y + 66), subtitle, fill=(219, 234, 254), font=font(20))
        if i < len(stages) - 1:
            ax = x + box_w + 7
            ay = y + box_h // 2
            draw.line((ax, ay, ax + gap - 14, ay), fill=(71, 85, 105), width=5)
            draw.polygon([(ax + gap - 14, ay - 12), (ax + gap + 4, ay), (ax + gap - 14, ay + 12)], fill=(71, 85, 105))
        x += box_w + gap
    draw.text((65, 85), "S.P.E.C.T.R.A Reconstruction Pipeline", fill=(15, 23, 42), font=font(56, True))
    draw.text((68, 158), "A dashboard-driven workflow from ordinary images to inspectable 3D assets", fill=(71, 85, 105), font=font(28))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def make_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets = {
        "pipeline": ASSET_DIR / "pipeline.png",
        "desk": ROOT / "data" / "results" / "fusion_quality_preview.png",
        "synthetic_car": ROOT / "data" / "results" / "synthetic" / "car_scene_preview.png",
        "synthetic_room": ROOT / "data" / "results" / "synthetic" / "indoor_room_dense_saturated_preview.png",
        "3drealcar": ASSET_DIR / "3drealcar_segmentation_panel.png",
        "hero": ASSET_DIR / "hero.png",
        "results": ASSET_DIR / "results.png",
    }
    make_pipeline_image(assets["pipeline"])
    if not assets["desk"].exists():
        make_placeholder(assets["desk"], "Desk Reconstruction", "Use latest dense cloud / mesh during live demo", (10, 14, 25), (37, 99, 235))
    if not assets["synthetic_car"].exists():
        make_placeholder(assets["synthetic_car"], "Synthetic Car", "Procedural fallback; 3DRealCar import path ready", (247, 249, 252), (37, 99, 235))
    if not assets["synthetic_room"].exists():
        make_placeholder(assets["synthetic_room"], "Synthetic Indoor Room", "RGB material-colored baseline", (247, 249, 252), (15, 118, 110))
    if not assets["3drealcar"].exists():
        make_placeholder(assets["3drealcar"], "3DRealCar Dataset", "Real car image and annotation assets", (247, 249, 252), (37, 99, 235))
    make_placeholder(assets["hero"], "S.P.E.C.T.R.A", "Smart Pipeline for 3D Reconstruction and Analysis", (15, 23, 42), (37, 99, 235))
    make_placeholder(assets["results"], "Recognizable Output", "Desk arrangement reconstructed from phone images", (15, 23, 42), (15, 118, 110))
    return assets


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.85))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.05), Inches(11.5), Inches(0.42))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLORS["muted"]


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(8)


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, color=COLORS["panel"]) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(226, 232, 240)
    tf = shape.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["ink"]
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLORS["muted"]


def add_image(slide, path: Path, x: float, y: float, w: float, h: float | None = None) -> None:
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def build_deck() -> None:
    assets = make_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    slide = prs.slides.add_slide(blank)
    add_image(slide, assets["hero"], 0, 0, 13.333, 7.5)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(5.45), Inches(10.8), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = "Image-Based 3D Reconstruction Dashboard"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    box2 = slide.shapes.add_textbox(Inches(0.92), Inches(6.05), Inches(10.8), Inches(0.45))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "Prepared for project presentation and live demonstration"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(203, 213, 225)

    # 2 Problem
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Problem Statement", "Turning ordinary phone images into usable 3D assets is still technically demanding.")
    add_bullets(slide, [
        "Manual photogrammetry workflows require many separate tools and careful file handling.",
        "Students and prototype teams need a visible end-to-end reconstruction workflow.",
        "Real captures contain blur, weak texture, changing exposure, reflective surfaces, and incomplete overlap.",
        "The project aims to make the reconstruction stages inspectable, explainable, and exportable."
    ], 0.85, 1.8, 6.2, 4.5)
    add_card(slide, "Core Idea", "A Streamlit dashboard that organizes image ingestion, frame selection, depth estimation, pose recovery, dense fusion, cleanup, meshing, and export.", 7.55, 1.85, 4.7, 2.0)
    add_card(slide, "Presentation Focus", "Show that the pipeline can reconstruct recognizable real objects and explain where limitations come from.", 7.55, 4.05, 4.7, 1.75)

    # 3 Objectives
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Project Objectives")
    for i, (title, body) in enumerate([
        ("Ingest Images", "Accept real smartphone image sequences and stage them for processing."),
        ("Select Keyframes", "Keep sharp and diverse frames so reconstruction uses stronger visual evidence."),
        ("Recover Geometry", "Estimate depth, camera poses, sparse points, dense clouds, and meshes."),
        ("Evaluate Outputs", "Inspect point clouds, meshes, metrics, and export-ready files."),
        ("Support Baselines", "Use synthetic scenes and research datasets for controlled comparisons."),
        ("Explain Results", "Make every stage understandable during demonstration and assessment."),
    ]):
        x = 0.7 + (i % 3) * 4.15
        y = 1.65 + (i // 3) * 2.15
        add_card(slide, title, body, x, y, 3.55, 1.45)

    # 4 Pipeline
    slide = prs.slides.add_slide(blank)
    add_title(slide, "System Pipeline", "Each output in the dashboard corresponds to one stage in this flow.")
    add_image(slide, assets["pipeline"], 0.55, 1.55, 12.2, 4.9)

    # 5 Implementation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Implementation Stack")
    add_card(slide, "Dashboard", "Streamlit interface for uploads, live status, previews, output inspection, and synthetic baselines.", 0.8, 1.55, 5.8, 1.35)
    add_card(slide, "Computer Vision", "OpenCV ORB features, matching, essential matrix pose recovery, sharpness scoring, and image preprocessing.", 0.8, 3.1, 5.8, 1.35)
    add_card(slide, "Depth and 3D", "Depth Anything V2, NumPy, Open3D point clouds, DBSCAN cleanup, Poisson surface reconstruction.", 0.8, 4.65, 5.8, 1.35)
    add_card(slide, "Exports", "PLY and OBJ assets for inspection in the dashboard, MeshLab, CloudCompare, Blender, or Gaussian Splatting workflows.", 7.0, 1.55, 5.1, 1.35)
    add_card(slide, "Recent Improvements", "Lobby keyframes increased to 60; object fusion made more permissive for thin/planar captures; synthetic car import path added.", 7.0, 3.1, 5.1, 1.95)

    # 6 Output stages
    slide = prs.slides.add_slide(blank)
    add_title(slide, "What The Output Types Mean")
    rows = [
        ("Sparse cloud", "Feature matches used to estimate camera alignment."),
        ("Pose-aligned dense cloud", "Main colored point cloud from depth maps and camera poses."),
        ("Legacy dense cloud", "Older/fallback dense output kept for comparison."),
        ("Cleaned dense cloud", "Noise-reduced cloud used before meshing."),
        ("Poisson mesh", "Continuous surface generated from the cleaned point cloud."),
        ("OBJ mesh", "Portable mesh format for external 3D software."),
    ]
    for i, (title, body) in enumerate(rows):
        add_card(slide, title, body, 0.75 + (i % 2) * 6.1, 1.35 + (i // 2) * 1.65, 5.55, 1.08)

    # 7 Real result
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Best Real Reconstruction: Desk Arrangement", "The strongest current output is recognizable as a physical workstation scene.")
    add_image(slide, assets["results"], 0.65, 1.35, 5.7, 4.65)
    add_bullets(slide, [
        "Recognizable elements: desk/table surface, laptop/screens, cups, and surrounding objects.",
        "Point cloud and mesh outputs show that the pipeline is reconstructing real object geometry.",
        "Remaining gaps and stretched areas are expected from monocular depth, reflective surfaces, and pose noise.",
        "This is the best demonstration artifact for the live presentation."
    ], 6.8, 1.65, 5.6, 4.6, 17)

    # 8 Synthetic baselines
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Synthetic Baselines", "Synthetic scenes help explain controlled geometry versus real capture uncertainty.")
    add_image(slide, assets["synthetic_room"], 0.75, 1.4, 5.65, 4.4)
    add_image(slide, assets["synthetic_car"], 6.9, 1.4, 5.65, 4.4)
    add_card(slide, "Note", "The current car preview is a procedural fallback. The importer is ready to replace it with 3DRealCar once a scanned car model is extracted.", 1.2, 6.05, 10.9, 0.75)

    # 9 Limitations
    slide = prs.slides.add_slide(blank)
    add_title(slide, "3DRealCar Dataset Integration", "The downloaded archive provides real car images and segmentation labels.")
    add_image(slide, assets["3drealcar"], 0.45, 1.2, 12.45, 5.8)

    # 10 Limitations
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Limitations And Risks")
    add_bullets(slide, [
        "Textureless walls, glossy screens, shadows, and repeated patterns can reduce pose accuracy.",
        "Monocular depth provides relative depth, so scale and alignment can drift.",
        "Thin planar objects may need permissive fusion settings to avoid over-filtering.",
        "Poisson meshing can fill gaps or stretch surfaces where the point cloud is incomplete.",
        "Synthetic colors are generated/material colors unless they come from a real scan or texture dataset."
    ], 0.85, 1.45, 7.2, 4.9, 18)
    add_card(slide, "How We Addressed This", "Increased keyframes for room captures, loosened object fusion, preserved thin structures, and added a dataset-backed car import path.", 8.45, 2.0, 3.9, 2.5)

    # 11 Evaluation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Evaluation Strategy")
    for i, (title, body) in enumerate([
        ("Visual Quality", "Can the viewer recognize the object/scene arrangement?"),
        ("Point Count", "Are enough points produced after fusion and cleanup?"),
        ("Stage Completion", "Did depth, SfM, fusion, cleanup, mesh, and export stages complete?"),
        ("Export Readiness", "Can PLY/OBJ outputs be opened in external 3D tools?"),
        ("Failure Analysis", "Do artifacts align with known capture or algorithm limitations?"),
        ("Baseline Comparison", "Compare real reconstructions with synthetic controlled scenes."),
    ]):
        add_card(slide, title, body, 0.75 + (i % 3) * 4.15, 1.5 + (i // 3) * 2.1, 3.55, 1.35)

    # 12 Demo script
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Live Demo Flow")
    add_bullets(slide, [
        "1. Open dashboard and show dataset ingestion.",
        "2. Show raw image inventory and preview dropdown.",
        "3. Show keyframe selection and depth preview outputs.",
        "4. Open 3D viewer and compare sparse cloud, dense cloud, cleaned cloud, and mesh.",
        "5. Explain the desk reconstruction as the best real output.",
        "6. Open synthetic baselines and mention 3DRealCar integration path.",
        "7. Close with limitations and future improvements."
    ], 1.0, 1.35, 11.0, 5.4, 20)

    # 13 Conclusion
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Conclusion")
    add_bullets(slide, [
        "S.P.E.C.T.R.A demonstrates a complete image-to-3D workflow in one dashboard.",
        "The project produces inspectable point clouds, meshes, and export-ready files.",
        "The best real result reconstructs a recognizable desk arrangement from phone images.",
        "The system is presentation-ready as a prototype, with clear explanations for artifacts and limitations.",
        "Future work: stronger calibration, real scanned car baseline, texture projection, and measured benchmark evaluation."
    ], 1.0, 1.45, 10.9, 4.8, 21)
    add_card(slide, "Final Message", "The project proves the workflow and makes every reconstruction stage visible, explainable, and demonstrable.", 1.0, 6.15, 11.0, 0.75, COLORS["panel"])

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build_deck()
